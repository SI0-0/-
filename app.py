import streamlit as st
from pptx import Presentation
from collections import defaultdict

# --------------------------
# 글꼴 다운로드 링크 매핑 (예시)
FONT_LINKS = {
    "Roboto": "https://fonts.google.com/specimen/Roboto",
    "Open Sans": "https://fonts.google.com/specimen/Open+Sans",
    "Arial": "https://www.wfonts.com/font/arial",
    "Times New Roman": "https://www.wfonts.com/font/times-new-roman",
    "Calibri": "https://www.wfonts.com/font/calibri",
    "Verdana": "https://www.wfonts.com/font/verdana",
    # 필요시 추가 가능
}
# --------------------------

# --------------------------
# Streamlit GUI 스타일
st.set_page_config(page_title="PPT 글꼴 확인기", page_icon="🎨", layout="centered")
st.markdown("""
<style>
.title { font-size:2.4rem; font-weight:700; text-align:center; margin-bottom:0.3rem; }
.subtitle { text-align:center; font-size:1.1rem; color:#666; margin-bottom:2rem; }
.card { padding:1rem; border-radius:16px; background:#f9f9f9; margin-bottom:1rem; border:1px solid #e3e3e3; box-shadow:0px 3px 10px rgba(0,0,0,0.05);}
.card a { text-decoration:none; color:#1f77b4; font-weight:600;}
</style>
""", unsafe_allow_html=True)

st.markdown('<div class="title">🎨 PPT 글꼴 확인기</div>', unsafe_allow_html=True)
st.markdown('<div class="subtitle">업로드한 PPTX 파일에서 사용된 글꼴과 스타일을 추출하고 다운로드 링크를 제공합니다.</div>', unsafe_allow_html=True)
# --------------------------

uploaded_file = st.file_uploader("📤 PPTX 파일 선택", type=["pptx"])

if uploaded_file:
    st.info("⚡ 분석 중...")

    prs = Presentation(uploaded_file)
    font_info = defaultdict(list)

    # 모든 슬라이드, 도형, 텍스트 분석
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        font_info[run.font.name].append({
                            "slide": slide_idx,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "size": run.font.size.pt if run.font.size else None
                        })

    if font_info:
        st.success(f"✅ 발견된 글꼴 {len(font_info)}개")

        # 카드형 GUI + 다운로드 링크 제공
        for font, details in font_info.items():
            link = FONT_LINKS.get(font)
            if link:
                st.markdown(
                    f'<div class="card"><b>{font}</b> - 사용 횟수: {len(details)} → '
                    f'<a href="{link}" target="_blank">다운로드 / 사이트 방문</a></div>', unsafe_allow_html=True)
            else:
                st.markdown(
                    f'<div class="card"><b>{font}</b> - 사용 횟수: {len(details)} → 링크 없음</div>', unsafe_allow_html=True)
    else:
        st.warning("❌ PPT에서 글꼴을 찾지 못했습니다.")
